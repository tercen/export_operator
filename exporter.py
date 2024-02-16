from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import datetime
from io import BytesIO
from hashlib import md5
from base64 import b64encode

import polars as pl
from os.path import basename
import subprocess
import random, string

class Exporter:
    def __init__(self, tmpFolder=None):
        self.pages = []

        if tmpFolder is None:
            self.tmpFolder = "/tmp/" + ''.join(random.choice(string.ascii_letters) for _ in range(15))
        else: 
            self.tmpFolder = tmpFolder


    def add_blank_page(self):
        return -1

    # Image info is a list [imgPath, imgSize]
    def add_image(self, imgInfo, page_idx=0):
        pass

    def add_title(self, title, page_idx=0):
        pass

    def add_text(self, textFile, page_idx=0, text_size=11):
        pass

    def get_bytes(self):
        return None
    
    def as_dataframe(self):
        return None


class PPTXExporter(Exporter):
    def __init__(self, output="pptx", tmpFolder=None):
        super().__init__(tmpFolder=tmpFolder)
        self.output = output
        self.presentation = Presentation()

    def add_blank_page(self, stpName=None):
        blank_slide_layout = self.presentation.slide_layouts[6]
        self.pages.append(self.presentation.slides.add_slide(blank_slide_layout))

        return len(self.pages)
    
    def add_title(self, title, page_idx=None):
        if page_idx == None:
            page_idx = len(self.pages)-1

        left = Inches(0)
        width = Inches(10)
        text_size = 16
        height = Inches(2)
        top = Inches(0.25)
        
        txBox = self.pages[page_idx].shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = title
        txBox.text_frame.paragraphs[0].font.size = Pt(text_size)
        
        txBox.text_frame.paragraphs[0].font.bold = True

    
    def add_footer(self, page_idx=None):
        if page_idx == None:
            page_idx = len(self.pages)-1

        now = datetime.datetime.now().strftime("%H:%M:%S --- %d, %B %Y")

        text_size = 10
        left = Inches(0)
        width = Inches(10)
        height = Inches(2)
        top = Inches(7.1)
        
        txBox = self.pages[page_idx].shapes.add_textbox(left, top, width, height)
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text =  now
        run.font.color.rgb = RGBColor(128, 128, 128)
        txBox.text_frame.paragraphs[0].font.size = Pt(text_size)


        top = Inches(7)
        left = Inches(8.48)
        height = Inches(0.5)

        self.pages[page_idx].shapes.add_picture("img/tercen.png", left, top, height=height)

    def add_image(self, imgInfo, page_idx=None):
        if page_idx is None:
            page_idx = len(self.pages)-1
        
        # Slide is 10 x 7.5 in. (w x h)
        top = Inches(0.75)
        left = Inches(0.5)
        height = Inches(5.8)
        #pic = 
        self.pages[page_idx].shapes.add_picture(imgInfo[0], left, top, height=height)

 

    def add_text(self, textFile, page_idx=0, text_size=11):

        with open(textFile, "r") as file:
            text = file.read()

        left = Inches(0)
        width = Inches(10)
        height = Inches(5.8)
        top = Inches(1.5)

        txBox = self.pages[page_idx].shapes.add_textbox(left, top, width, height)

        txBox.text_frame.text = text
        txBox.text_frame.paragraphs[0].font.size = Pt(text_size)




    def get_bytes(self):
        buf = BytesIO()
        self.presentation.save(buf)
        return buf
    
    def as_dataframe(self, filename):
        
        self.presentation.save(self.tmpFolder + "/" + basename(filename) + ".pptx")
        
        

        if self.output == "pptx":
            outname = basename(filename) + ".pptx"
            mimetype = "application/vnd.ms-powerpoint"
            with open(self.tmpFolder + "/" + basename(filename) + ".pptx", "rb") as file:
                fileBytes = file.read()
        else:
            outname = basename(filename) + ".pdf"
            mimetype = "application/pdf"
            subprocess.call(["libreoffice", "--headless" ,\
                            "--convert-to", "pdf", \
                            self.tmpFolder + "/" + basename(filename) + ".pptx", \
                            "--outdir", self.tmpFolder])
            with open(self.tmpFolder + "/" + basename(filename) + ".pdf", "rb") as file:
                fileBytes = file.read()
        
        checksum = md5(fileBytes).hexdigest()

        imgDf = pl.DataFrame({\
            ".ci":[0],\
            "filename":[outname],\
            "mimetype":[mimetype],\
            "checksum":[checksum],\
            ".content":[str(b64encode(fileBytes).decode("utf-8"))]\
            })
        
        imgDf = pl.with_column(pl.col('bar').cast(pl.Int32))
        return imgDf
    

# def image_file_to_df(file_path):
#     checksum = hashlib.md5(open(file_path,'rb').read()).hexdigest()

#     output_str = []

#     for fpath in file_path:
#         with open(file_path, mode="rb") as f:
#             fc = f.read()
#             output_str.append([base64.b64encode(fc)])


#     o = output_str[0][0]

#     outs = o.decode('utf-8')
#     imgDf = pd.DataFrame({
#         "filename":[filename],
#         "mimetype":[mimetype],
#         "checksum":[checksum],
#         ".content":[outs]
#     })

#     return imgDf