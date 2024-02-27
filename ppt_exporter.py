from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import datetime
from hashlib import md5
from base64 import b64encode

import polars as pl
from os.path import basename
import subprocess

import re



from PIL import Image
from exporter import Exporter


class PPTXExporter(Exporter):
    def __init__(self, output="pptx", tmpFolder=None):
        super().__init__(tmpFolder=tmpFolder)
        self.output = output
        self.presentation = Presentation()

    def add_blank_page(self, stpName=None):
        blank_slide_layout = self.presentation.slide_layouts[6]
        self.pages.append(self.presentation.slides.add_slide(blank_slide_layout))

        return len(self.pages)
    
    def add_title(self, title, page_idx=None):
        if page_idx == None:
            page_idx = len(self.pages)-1

        left = Inches(0)
        width = Inches(10)
        text_size = 16
        height = Inches(2)
        top = Inches(0.25)
        
        txBox = self.pages[page_idx].shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = title
        txBox.text_frame.paragraphs[0].font.size = Pt(text_size)
        
        txBox.text_frame.paragraphs[0].font.bold = True

    
    def add_footer(self, page_idx=None):
        if page_idx == None:
            page_idx = len(self.pages)-1

        now = datetime.datetime.now().strftime("%H:%M:%S --- %d, %B %Y")

        text_size = 10
        left = Inches(0)
        width = Inches(10)
        height = Inches(2)
        top = Inches(7.1)
        
        txBox = self.pages[page_idx].shapes.add_textbox(left, top, width, height)
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text =  now
        run.font.color.rgb = RGBColor(128, 128, 128)
        txBox.text_frame.paragraphs[0].font.size = Pt(text_size)


        top = Inches(7)
        left = Inches(8.48)
        height = Inches(0.5)

        self.pages[page_idx].shapes.add_picture("img/tercen.png", left, top, height=height)

    def finish_page(self):
        pass

    def add_image(self, imgInfo, page_idx=None):
        if page_idx is None:
            page_idx = len(self.pages)-1
        
        # Slide is 10 x 7.5 in. (w x h)
        im = Image.open(imgInfo[0])
        width, height = im.size


        top = Inches(0.75)
        left = Inches(0.5)
       
        pg = self.pages[page_idx]
        
        pgImg = pg.shapes.add_picture(imgInfo[0], left, top, height=Inches(5.8))

        
        heightRel = pgImg.height / self.presentation.slide_height
        widthRel = pgImg.width / self.presentation.slide_width


        if heightRel > 0.8 or widthRel > 0.8:
            if widthRel > heightRel:
                relChange = 0.8 / widthRel
                pgImg.width = int(pgImg.width * relChange * 0.8)
                pgImg.height = int(pgImg.height * relChange * 0.8)
            else:
                relChange = 0.8 / heightRel
                pgImg.width = int(pgImg.width * relChange * 0.8)
                pgImg.height = int(pgImg.height * relChange * 0.8)
            

    def add_text(self, textFile, page_idx=0, text_size=11):

        with open(textFile, "r") as file:
            text = file.readlines()

        line=text[0]
        y = 1.0
        left = Inches(0)
        width = Inches(5)
        height = Inches(5.8)
        top = Inches(1)
        txBox = self.pages[page_idx].shapes.add_textbox(left, top, width, height)


        tableLines = []
        tableStart = None
        drawingTable = False
        for line in text:
            y += 0.19
            if line.startswith("#"):
                run = txBox.text_frame.paragraphs[0].add_run()
                run.font.bold = True
                
                if line.startswith("###"):
                    run.text =  line.split("### ")[1]
                    run.font.size = Pt(text_size + 2)
                    
                if line.startswith("###"):
                    run.text =  line.split("## ")[1]
                    run.font.size = Pt(text_size + 4)

                if line.startswith("###"):
                    run.text =  line.split("# ")[1]
                    run.font.size = Pt(text_size + 6)
            elif line.startswith("|"):
                tableLines.append(line)

                if drawingTable == False:
                    tableStart = y 
                    drawingTable = True

                # Add blank spaces in case more text comes after the table
                run = txBox.text_frame.paragraphs[0].add_run()
                run.text =  "\n"
                run.font.size = Pt(text_size+1)
                run.font.bold = False
            else:

                run = txBox.text_frame.paragraphs[0].add_run()
                run.text =  line
                run.font.size = Pt(text_size)
                run.font.bold = False

        

        if len(tableLines) > 0:
            header = tableLines[0]

            nCols = len(re.findall("[|]", header)) - 1
            nRows = len(tableLines) - sum([l.startswith("|-") and l.endswith("-|\n") for l in tableLines])

            #TODO, check for header marking
            shape = self.pages[page_idx].shapes.add_table(nRows, nCols, \
                                            Inches(0.8), 
                                            Inches(tableStart),
                                            Inches(5),
                                            Inches( len(tableLines) * 0.25  ))
            
            tbl = shape.table
            
            ri = 0
            for line in tableLines:
                if line.startswith("|-") and ":|:" in line:
                    continue

                content = line.split("|")
                ci = 0
                for i in range(1, len(content)-1):
                    if ri == 0:
                        tbl.cell(ri,ci).fill.solid()
                        tbl.cell(ri,ci).fill.fore_color.rgb = RGBColor(128, 128, 128)
                    else:
                        tbl.cell(ri,ci).fill.solid()
                        tbl.cell(ri,ci).fill.fore_color.rgb = RGBColor(220, 220, 220)

                    run = tbl.cell(ri, ci).text_frame.paragraphs[0].add_run()
                    run.text =  str.strip(content[i])
                    
                    if ri == 0:
                        run.font.size = Pt(text_size+3)
                        run.font.bold = True
                    else:
                        run.font.size = Pt(text_size)
                        run.font.bold = False

                    ci += 1

                ri += 1

   
    def as_dataframe(self, filename):
        # self.presentation.save("test.pptx")
        self.presentation.save(self.tmpFolder + "/" + basename(filename) + ".pptx")
        
        

        if self.output == "pptx":
            outname = basename(filename) + ".pptx"
            mimetype = "application/vnd.ms-powerpoint"
            with open(self.tmpFolder + "/" + basename(filename) + ".pptx", "rb") as file:
                fileBytes = file.read()
        else:
            outname = basename(filename) + ".pdf"
            mimetype = "application/pdf"
            subprocess.call(["libreoffice", "--headless" ,\
                            "--convert-to", "pdf", \
                            self.tmpFolder + "/" + basename(filename) + ".pptx", \
                            "--outdir", self.tmpFolder])
            with open(self.tmpFolder + "/" + basename(filename) + ".pdf", "rb") as file:
                fileBytes = file.read()
        
        checksum = md5(fileBytes).hexdigest()

        imgDf = pl.DataFrame({\
            ".ci":[0],\
            "filename":[outname],\
            "mimetype":[mimetype],\
            "checksum":[checksum],\
            ".content":[str(b64encode(fileBytes).decode("utf-8"))]\
            })
        
        imgDf = imgDf.with_columns(pl.col('.ci').cast(pl.Int32))
        return imgDf
